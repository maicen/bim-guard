"""
build_fmp_decks.py
------------------------------------------------
Builds the four BIMGUARD AI FMP presentations (96 slides) in the visual
style of docs/Presentation-2026-May/BIMGUARD_AI_FMP_Presentation.pptx.

STYLE
    Extracted from the template rather than approximated: 10.00 x 5.62in
    slides, Calibri throughout, navy #1F3864 / cyan #00B4D8 / blue #2E75B6
    palette, a 0.12in left accent bar, 26pt bold navy titles over a 13pt
    grey subtitle and a #DEEAF1 rule, and a navy footer band carrying the
    same string on every slide. Tables are drawn as shape grids because the
    template contains no native PowerPoint tables.

NUMBERS
    Every measured figure comes from this repository's own outputs —
    validation_sweep_summary.json, docs/validation/*.csv, and the pytest
    runs — not from the drafting brief. Where the brief specified figures
    that are not computed anywhere (per-engine precision and recall,
    ARCH rule accuracy, 2-8s runtimes), the measured value is shown
    instead and the slide says what was actually observed.

    That substitution is deliberate. Thesis §13.9 states that the MM-001
    and XM-001 packs are not wired into the comparator and that "no claim
    in this chapter should be read as asserting that it has"; slides
    quoting precision figures for engines that never execute would
    contradict the thesis they accompany. The engine-status slides
    therefore report coverage and flag rates, which are measurable, and
    name §13.9 as the reference.

    Presentation C is titled as a specification and design proposal
    because no ARCH engine exists in the codebase: a search for AR-001
    through AR-006 returns nothing. Its content is preserved as the
    designed ruleset, marked as not yet implemented.

Usage:
    uv run python build_fmp_decks.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT_DIR = Path("docs/Presentation-2026-May")
TEMPLATE = OUT_DIR / "BIMGUARD_AI_FMP_Presentation.pptx"

# ── Palette, lifted from the template ──────────────────────────────────────
NAVY = RGBColor(0x1F, 0x38, 0x64)
CYAN = RGBColor(0x00, 0xB4, 0xD8)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
RED = RGBColor(0xA3, 0x2D, 0x2D)
AMBER = RGBColor(0x85, 0x4F, 0x0B)
GREEN = RGBColor(0x1E, 0x7A, 0x4B)
GREY = RGBColor(0x88, 0x96, 0xA6)
BODY = RGBColor(0x37, 0x41, 0x51)
RULE = RGBColor(0xDE, 0xEA, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SHADE = RGBColor(0xF3, 0xF6, 0xF9)

FONT = "Calibri"
FOOTER = ("BIMGUARD AI  |  Final Master Project  |  "
          "Zigurat Global Institute of Technology  |  Group 5")

BAND = [CYAN, BLUE, NAVY, AMBER, RED, GREEN]


# ── Primitives ─────────────────────────────────────────────────────────────


def _box(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def _text(slide, x, y, w, h, runs, size=11, colour=BODY, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """runs: a string, or a list of (text, size, colour, bold) tuples."""
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(item, str):
            t, sz, col, bd = item, size, colour, bold
        else:
            t, sz, col, bd = item
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = col
    return sh


def _footer(slide):
    _box(slide, 0, 5.35, 10.0, 0.28, fill=NAVY)
    _text(slide, 0.25, 5.36, 9.5, 0.25, FOOTER, size=8.5, colour=WHITE)


def title_slide(prs, title, subtitle, standards, focus):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    _box(s, 0, 0, 10.0, 5.62, fill=NAVY)
    _box(s, 0, 0, 0.18, 5.62, fill=CYAN)
    _text(s, 0.40, 0.90, 9.2, 1.2, title, size=54, colour=WHITE, bold=True)
    _text(s, 0.40, 2.10, 9.2, 0.55, subtitle, size=22, colour=CYAN)
    _text(s, 0.40, 2.75, 9.2, 0.40, standards, size=13, colour=GREY)
    _box(s, 0.40, 3.22, 4.5, 0.02, fill=CYAN)
    _text(s, 0.40, 3.40, 9.2, 0.60, focus, size=12, colour=WHITE)
    _text(s, 0.25, 4.83, 9.5, 0.25, FOOTER, size=9, colour=GREY)
    return s


def slide(prs, title, subtitle, accent=CYAN):
    """A content slide: accent bar, title, subtitle, rule, footer."""
    s = prs.slides.add_slide(prs.slide_layouts[0])
    _box(s, 0, 0, 0.12, 5.62, fill=accent)
    _text(s, 0.30, 0.28, 9.4, 0.55, title, size=26, colour=NAVY, bold=True)
    _text(s, 0.30, 0.82, 9.4, 0.32, subtitle, size=13, colour=GREY)
    _box(s, 0.30, 1.18, 9.4, 0.02, fill=RULE)
    _footer(s)
    return s


def metrics(s, items, y=1.35, h=1.50):
    """Up to 4 metric boxes: (value, caption, colour)."""
    n = len(items)
    w, gap, x0 = 2.20, 0.20, 0.25
    if n < 4:
        w = (9.5 - gap * (n - 1)) / n
    for i, (val, cap, col) in enumerate(items):
        x = x0 + i * (w + gap)
        _box(s, x, y, w, h, fill=WHITE, line=RULE)
        _box(s, x, y, w, 0.08, fill=col)
        _text(s, x + 0.10, y + 0.18, w - 0.20, 0.78, val, size=36, colour=col, bold=True)
        _text(s, x + 0.10, y + 0.96, w - 0.20, h - 1.02, cap, size=11, colour=GREY, spacing=0.95)


def cards(s, items, y=1.35, h=None, gap=0.15):
    """Stacked full-width cards: (heading, body, colour)."""
    n = len(items)
    if h is None:
        h = min(1.12, (3.85 - gap * (n - 1)) / n)
    for i, (head, body, col) in enumerate(items):
        yy = y + i * (h + gap)
        _box(s, 0.25, yy, 9.5, h, fill=WHITE, line=RULE)
        _box(s, 0.25, yy, 0.07, h, fill=col)
        _text(s, 0.39, yy + 0.10, 9.28, 0.30, head, size=13, colour=NAVY, bold=True)
        _text(s, 0.39, yy + 0.40, 9.28, h - 0.48, body, size=11, colour=BODY, spacing=0.95)


def columns(s, items, y=1.35, h=3.30):
    """2-4 side-by-side boxes: (heading, body, colour)."""
    n = len(items)
    gap = 0.18
    w = (9.5 - gap * (n - 1)) / n
    for i, (head, body, col) in enumerate(items):
        x = 0.25 + i * (w + gap)
        _box(s, x, y, w, h, fill=WHITE, line=RULE)
        _box(s, x, y, w, 0.07, fill=col)
        _text(s, x + 0.12, y + 0.20, w - 0.24, 0.46, head, size=12, colour=NAVY, bold=True)
        _text(s, x + 0.12, y + 0.68, w - 0.24, h - 0.80, body, size=10.5, colour=BODY, spacing=0.95)


def table(s, headers, rows, y=1.35, widths=None, size=11, row_h=0.30):
    """Shape-grid table matching the template (no native PPT tables)."""
    total = 9.5
    widths = widths or [total / len(headers)] * len(headers)
    scale = total / sum(widths)
    widths = [w * scale for w in widths]
    _box(s, 0.25, y, total, row_h, fill=NAVY)
    x = 0.25
    for w, htxt in zip(widths, headers):
        _text(s, x + 0.07, y + 0.05, w - 0.14, row_h - 0.08, htxt, size=size - 0.5,
              colour=WHITE, bold=True)
        x += w
    for ri, row in enumerate(rows):
        yy = y + row_h + ri * row_h
        _box(s, 0.25, yy, total, row_h, fill=SHADE if ri % 2 else WHITE, line=RULE)
        x = 0.25
        for w, cell in zip(widths, row):
            txt, col, bold = (cell if isinstance(cell, tuple) else (cell, BODY, False))
            _text(s, x + 0.07, yy + 0.05, w - 0.14, row_h - 0.08, str(txt), size=size,
                  colour=col, bold=bold)
            x += w
    return y + row_h * (len(rows) + 1)


def formula(s, text, y=1.45, note=None, h=0.80):
    _box(s, 0.25, y, 9.5, h, fill=SHADE, line=RULE)
    _text(s, 0.35, y + 0.18, 9.3, h - 0.3, text, size=18, colour=NAVY, bold=True,
          align=PP_ALIGN.CENTER)
    if note:
        _text(s, 0.25, y + h + 0.10, 9.5, 0.30, note, size=10.5, colour=GREY,
              align=PP_ALIGN.CENTER)


def note(s, text, y=4.55, bold=False, colour=None):
    _text(s, 0.25, y, 9.5, 0.70, text, size=12 if bold else 11.5,
          colour=colour or (NAVY if bold else BODY), bold=bold, spacing=0.95)


def closing(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    _box(s, 0, 0, 10.0, 5.62, fill=NAVY)
    _box(s, 0, 0, 0.18, 5.62, fill=CYAN)
    _text(s, 0.40, 1.90, 9.2, 0.90, title, size=44, colour=WHITE, bold=True)
    _box(s, 0.40, 2.90, 4.5, 0.02, fill=CYAN)
    _text(s, 0.40, 3.10, 9.2, 0.50, subtitle, size=22, colour=CYAN)
    _text(s, 0.25, 4.83, 9.5, 0.25, FOOTER, size=9, colour=GREY)
    return s


def new_deck() -> Presentation:
    prs = Presentation(str(TEMPLATE))
    for i in range(len(prs.slides) - 1, -1, -1):  # strip template slides, keep master
        rid = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[i]
    return prs
